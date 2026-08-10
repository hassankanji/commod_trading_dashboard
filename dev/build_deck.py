"""Builds the intern presentation on the Standard Chartered Prosper template.

Only the template's own layouts and placeholders are used, so the design is the
template's rather than something bolted on top. Charts go in at the content
placeholder's own geometry.

    SC_TEMPLATE=/path/to/template.pptx python dev/build_deck.py
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import PP_PLACEHOLDER

HERE = os.path.dirname(os.path.abspath(__file__))
TEMPLATE = os.environ.get(
    "SC_TEMPLATE",
    os.path.join(HERE, "template_standard_chartered.pptx"))
OUT = os.environ.get("DECK_OUT",
                     os.path.join(os.path.dirname(HERE),
                                  "commodities_dashboard_presentation.pptx"))
ASSETS = os.environ.get("DECK_ASSETS", os.path.join(HERE, "assets"))

CHROME = (PP_PLACEHOLDER.DATE, PP_PLACEHOLDER.FOOTER)

prs = Presentation(TEMPLATE)
M0, M1 = prs.slide_masters[0], prs.slide_masters[1]

ids = prs.slides._sldIdLst
for sld in list(ids):
    prs.part.drop_rel(sld.rId)
    ids.remove(sld)


def add(layout):
    s = prs.slides.add_slide(layout)
    for ph in list(s.placeholders):
        if ph.placeholder_format.type in CHROME:
            ph._element.getparent().remove(ph._element)
    return s


def drop(shape):
    shape._element.getparent().remove(shape._element)


def title(slide, text):
    slide.placeholders[0].text_frame.paragraphs[0].add_run().text = text


def bullets(ph, items, size=None):
    tf = ph.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        bold, text = (it if isinstance(it, tuple) else (False, it))
        run = para.add_run()
        run.text = text
        run.font.bold = bold
        if size:
            run.font.size = Pt(size)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


def picture(slide, ph, path):
    from PIL import Image
    left, top, cw, ch = ph.left, ph.top, ph.width, ph.height
    drop(ph)
    iw, ih = Image.open(path).size
    scale = min(cw / iw, ch / ih)
    w, h = int(iw * scale), int(ih * scale)
    slide.shapes.add_picture(path, left + (cw - w) // 2, top + (ch - h) // 2, w, h)


# ---------------------------------------------------------------- 1. Title
s = add(M0.slide_layouts[0])
title(s, "Commodities Vol & Relative-Value Dashboard")
s.placeholders[1].text_frame.paragraphs[0].add_run().text = \
    "Hassan Kanji  |  Summer Internship 2026"
drop(s.placeholders[10])
notes(s, """
Ten minutes. What I built, how it decides something is worth trading, and whether those calls
were actually right. Most of my time went into that last part, and it changed what I think of
the tool.
""")

# ------------------------------------------------------- 2. What I built
s = add(M1.slide_layouts[28])
title(s, "What I built")
bullets(s.placeholders[2], [
    "A daily monitor over 29 commodity and macro instruments. It looks for places where "
    "volatility or a relationship has moved outside its normal range, and writes each one up as "
    "a trade: what to do, what triggered it, and what has to happen for it to pay."])
notes(s, """
Screenshot of the Signals tab goes in the picture placeholder.

Six commodity groups: energy, precious metals, base metals, grains, softs, livestock. Plus S&P,
the dollar and ten-year futures, so commodities can be tested against the things that push them
around. 23 of the 29 have listed options, so those carry implied volatility as well as price.

It does not forecast prices. It notices when something is priced oddly against its own history
and says so as a trade rather than as a statistic.
""")

# ------------------------------------------------- 3. The shared idea
s = add(M1.slide_layouts[8])
title(s, "Every engine asks the same question")
bullets(s.placeholders[1], [
    (True, "Is this number unusual for this asset, by its own history?"),
    "Each quantity gets ranked against its own past year. Top or bottom 10% counts as stretched.",
    (True, "Example: WTI implied vol at the 8th percentile"),
    "Options on WTI are pricing less risk than this market has priced on 92% of the past year's "
    "days. That says nothing about where oil is going. It says volatility looks cheap against "
    "how this market usually prices it, so the trade is to buy volatility.",
    (True, "The bet is mean reversion"),
    "Stretched things snap back more often than they stretch further. Every engine runs that "
    "same test on a different quantity, and the backtest later is a test of that assumption.",
], size=16)
notes(s, """
Slow down here. If they take one thing away, it is that the tool measures a quantity against its
own history and calls the extremes.

Why percentiles instead of levels: 20 vol is cheap for natural gas and expensive for gold.
Ranking each asset against itself makes one threshold work across everything.

If asked why 10%: it is a slider, not a constant. Loosen it and you get more ideas of lower
quality. The backtest measures that trade-off.
""")

# --------------------------------------------------- 4. Five engines
s = add(M1.slide_layouts[8])
title(s, "Five engines, five kinds of dislocation")
ph = s.placeholders[1]
rows = [("Engine", "What it looks for", "The trade"),
        ("Vol dispersion",
         "Two related assets whose vol spread has stretched from its norm",
         "Sell the expensive vol, buy the cheap one"),
        ("Variance risk premium",
         "Implied vol well above the volatility actually being delivered",
         "Sell vol and collect the premium"),
        ("IV mean-reversion",
         "Implied vol at the top or bottom of its own 1-year range",
         "Buy vol when cheap, sell it when rich"),
        ("Correlation RV",
         "Two assets that normally track, now pulled apart",
         "Buy the laggard, sell the outperformer"),
        ("Lead-lag catch-up",
         "One asset that reliably moves a few days before another",
         "Trade the follower in the leader's direction")]
left, top, width, height = ph.left, ph.top, ph.width, ph.height
drop(ph)
tbl = s.shapes.add_table(len(rows), 3, left, top, width, height).table
tbl.columns[0].width = Inches(3.0)
tbl.columns[1].width = Inches(5.3)
tbl.columns[2].width = Inches(4.36)
for r, row in enumerate(rows):
    for c, val in enumerate(row):
        run = tbl.cell(r, c).text_frame.paragraphs[0].add_run()
        run.text = val
        run.font.size = Pt(14 if r == 0 else 13)
        run.font.bold = bool(r == 0 or c == 0)
notes(s, """
Do not read the table. Explain two properly and say the other three follow the same pattern.

Take the variance risk premium, because that is where a lot of the edge turns out to be. Option
markets tend to charge more for volatility than the asset ends up delivering. When that gap gets
unusually wide, you sell vol and collect it.

And take correlation RV, because people misread it. It is a convergence trade. When two names
that normally move together pull apart, it buys the one that lagged and sells the one that ran.
Lead-lag is the opposite: it follows the leader.

They are listed here in the order the backtest ranked them, best first. The first three trade
volatility so they are measured in vol points. The last two trade price, measured in percent.
That is why the results table has two different units.
""")

# ------------------------------------------ 5. Supporting views
s = add(M1.slide_layouts[31])
title(s, "Two supporting views")
bullets(s.placeholders[2], [
    "Vol Monitor ranks every asset by implied vol, realized vol and the gap between them, with a "
    "data-quality flag so a stale option mark cannot pass as a signal.\n\n"
    "Correlation & Pairs shows how tightly two assets move together, whether that link is "
    "stretched right now, and which of the two moves first."])
notes(s, """
Screenshot of the Correlation & Pairs tab goes in the picture placeholder.

The data-quality flag is computed, not assigned by hand. It measures how often an implied-vol
series actually re-prices. Even liquid commodity options only print a fresh mark on 50 to 70% of
days on this feed, and thin ones barely move. Signals built on a frozen mark get excluded.

That flag turns out to matter more than I expected, which comes up two slides later.

The lead-lag chart is the one to point at. It shows correlation at every lag, so a flat curve
means there is no timing edge, just ordinary correlation.
""")

# --------------------------------------------- 6. The test
s = add(M1.slide_layouts[12])
title(s, "How I tested whether any of it works")
bullets(s.placeholders[1], [
    (True, "Replay a year, one week at a time"),
    "52 replays. At each date the data is rewound and the dashboard's own code runs again, so "
    "the signals being graded are the ones it would really have printed. It cannot see past the "
    "replay date.",
    (True, "Grade every signal, all 1,970 of them"),
    "Not a hand-picked few. One unit per trade, no sizing, marked from the close on the signal "
    "date to the close five days later.",
    (True, "Compare each one against doing nothing"),
    "A win rate on its own proves nothing. So each trade is also scored as if it had been taken "
    "on every date in the year. That is the baseline, and the gap to it is the only number I "
    "would call an edge.",
], size=15)
picture(s, s.placeholders[13], "%s/hold.png" % ASSETS)
notes(s, """
Sell the baseline here, because it is what makes the rest mean anything.

Concrete version: selling volatility won about 59% of the time. Sounds great, until you check
that selling volatility on a random day with no signal won about 50% of the time, because
implied vol usually runs above realized. The signal is worth the gap, not the 59%.

Without that comparison I would have reported five working engines. With it, two disappear.

The chart answers the obvious question about why five days. Five is not the best number. The
edge more than triples if you hold to 21 days. But replays are five days apart, so five-day
holds are the only ones that do not overlap each other. Anything longer double-counts the same
market moves and makes the p-values look better than they are. So I quote the conservative one
and show the rest.

The shape is a finding too. Almost nothing happens in the first three days. These dislocations
take weeks to resolve, so a one-week hold is probably too short if you were actually trading it.
""")

# --------------------------------------------- 7. Results
s = add(M1.slide_layouts[8])
title(s, "Three of the five engines beat doing nothing")
picture(s, s.placeholders[1], "%s/engines.png" % ASSETS)
notes(s, """
Across all 1,970 trades the signals won 54.3% against a 48.7% baseline. A 5.6 point edge, 95%
interval 52.1 to 56.5, p = 0.0000. Real, and unevenly spread.

Vol dispersion is the standout: 64.1% against 48.2%, so nearly 16 points, on 192 trades, and
+12.1 vol points per trade. The variance risk premium is +8.9 and IV mean-reversion +7.7, both
significant.

The two price engines do not work. Lead-lag is +1.2 and has the largest sample of the lot at 598
trades, which is the useful way to put it: this is not a small-sample problem, it is the
most-tested engine and the edge still is not there. Correlation RV is +3.3 with p = 0.07, so
suggestive at best.

Worth explaining IV mean-reversion. It won 52.5%, barely a coin flip, which looks weak. But its
baseline is 44.8%, the lowest of any engine, because buying volatility loses more often than it
wins by nature. Vol drifts down most of the time and pays off in bursts. Against the right
benchmark that 52.5% is a genuine +7.7. Best example on the slide of why the baseline matters.
""")

# --------------------------------------------- 8. Where the edge lives
s = add(M1.slide_layouts[8])
title(s, "The edge is in agriculture, and in liquid options")
picture(s, s.placeholders[1], "%s/edge_lives.png" % ASSETS)
notes(s, """
Two cuts of the same trades.

By underlying, cocoa, the S&P and soybean meal come out best, all around +18 to +21 points. US
ten-year futures, natural gas and silver come out worst. By asset class the pattern is cleaner:
grains +9.3, base metals +9.1, softs +8.6, against energy +3.7 and macro +2.7. The edge sits in
agriculture and away from the most heavily traded markets.

Caveat to say out loud: 26 assets were tested, so a few will look significant by luck alone. WTI
is +9.0 and Brent is +1.0, and those two are nearly the same market. I would not trust any single
name. The class-level pattern is the reliable part.

The right panel is the one I did not expect, and it points the other way from what I assumed.
Signals built on options that re-price most days ran +11.9 points. Signals on options that
re-price rarely ran +6.2. I had guessed the thin markets would be less efficient and give a
better signal. The opposite is true, and the reason makes sense: if the option mark barely moves,
an apparent dislocation is partly stale data rather than a real mispricing. The data-quality
filter is not hygiene, it is where the edge is.

The price engines, which never touch options, sit at +2.5. Same story from another angle: the
volatility engines are carrying this.
""")

# --------------------------------------------- 9. Confidence
s = add(M1.slide_layouts[8])
title(s, "The confidence score does not pick winners")
picture(s, s.placeholders[1], "%s/confidence.png" % ASSETS)
notes(s, """
The dashboard puts a confidence percentage on every idea. It comes from the asset's own history:
on past days when this setup was at least this extreme, how often did the bet work.

It does not survive the test. Inside an engine, the rank correlation between the score and
whether the trade worked is -0.002. I checked every holding period from one day to 21 and it
never gets above +0.04. Split each engine at its own median and the less confident half actually
did better, +6.7 against +4.3.

Raising the minimum-confidence filter looks like it helps, but that is engine selection in
disguise. At a 70% minimum the edge is +7.2 points, and +8.7 of that is explained purely by which
engines survive the filter. Confidence itself contributes -1.5. You can get the same result by
switching off the two engines that do not work, without throwing away 84% of the sample.

Why it fails is the interesting part. The score measures how often a setup worked, not how much
more often it worked than the same bet on any other day. So an engine sitting on a favourable
base rate scores high with no skill in it. It is measuring the base rate and calling it
confidence.

The fix is not to delete it. Score it as the gap to the base rate instead of the raw hit rate.
I have a diagnostic version in the backtest; a deployable one needs a point-in-time base rate
rather than one computed over the whole sample.
""")

# --------------------------------------------- 10. Change and limits
s = add(M1.slide_layouts[12])
title(s, "A win rate means nothing until you know the baseline")
bullets(s.placeholders[1], [
    (True, "What I would change"),
    "Rebuild confidence as the gap to a base rate, not a raw hit rate.",
    "Drop correlation RV and lead-lag, or rework them. 1,104 trades between them and no edge "
    "worth the name.",
    "Lean on the volatility engines and on liquid option markets, which is where the edge sits.",
    "Show every signal against its baseline in the dashboard, not only in the backtest.",
], size=15)
bullets(s.placeholders[13], [
    (True, "What this does not prove"),
    "No costs, no sizing, no option greeks. This measures direction, not profit.",
    "A setup that stays extreme for weeks gets recorded at every replay, so the trades are not "
    "independent and the real sample is smaller than 1,970.",
    "One year, one universe, one set of parameters.",
    "26 assets tested, so some of the per-asset results are luck.",
], size=15)
notes(s, """
Land the title first. The dashboard was the build. The baseline was the judgement. Two engines
and the confidence score only looked good until there was something to compare them against.

Then say the limitations before anyone asks. That is the difference between a result and a claim.

The overlap point is the one a sharp listener raises. Weekly replays with five-day holds do not
overlap in time, but the same signal reappearing for six weeks is six correlated records, not six
independent ones. That inflates confidence in the result, which is why I would describe +5.6
points as small and real rather than proven.

If asked what I would do with more time: run the test across several years to see whether vol
dispersion holds up outside this particular one, and test the rebuilt confidence score properly
out of sample.

And if asked what I learned: a win rate without a benchmark is not evidence, building the thing
that tests your work matters as much as the work, and the most useful result was the one telling
me part of it did not work.
""")

prs.save(OUT)
print("wrote %s - %d slides" % (OUT, len(prs.slides)))
